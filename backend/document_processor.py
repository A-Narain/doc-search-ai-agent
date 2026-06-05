import fitz
from docx import Document
import os


# ── EXISTING: PDF ────────────────────────────────────────────────────────────
def extract_pdf_text(filepath):
    text = ""
    pdf = fitz.open(filepath)
    for page in pdf:
        text += page.get_text()
    return text


# ── EXISTING: DOCX ───────────────────────────────────────────────────────────
def extract_docx_text(filepath):
    doc = Document(filepath)
    return "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
    )


# ── NEW: Excel (.xlsx / .xls) ────────────────────────────────────────────────
def extract_excel_text(filepath):
    import openpyxl
    wb = openpyxl.load_workbook(filepath, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(
                str(c) for c in row if c is not None
            )
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


# ── NEW: PowerPoint (.pptx) ──────────────────────────────────────────────────
def extract_pptx_text(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"[Slide {i + 1}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


# ── NEW: Outlook Message (.msg) ──────────────────────────────────────────────
def extract_msg_text(filepath):
    import extract_msg
    msg = extract_msg.Message(filepath)
    return (
        f"From: {msg.sender}\n"
        f"To: {msg.to}\n"
        f"Subject: {msg.subject}\n"
        f"Date: {msg.date}\n\n"
        f"{msg.body or ''}"
    )


# ── NEW: HTML (.html / .htm) ─────────────────────────────────────────────────
def extract_html_text(filepath):
    from bs4 import BeautifulSoup
    with open(filepath, encoding="utf-8") as f:
        soup = BeautifulSoup(f, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


# ── NEW: Markdown / Plain Text (.md / .txt) ──────────────────────────────────
def extract_plain_text(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


# ── ROUTER: maps extension → handler ────────────────────────────────────────
def extract_text(filepath):
    extension = os.path.splitext(filepath)[1].lower()

    handlers = {
        ".pdf":  extract_pdf_text,
        ".docx": extract_docx_text,
        ".xlsx": extract_excel_text,
        ".xls":  extract_excel_text,
        ".pptx": extract_pptx_text,
        ".msg":  extract_msg_text,
        ".html": extract_html_text,
        ".htm":  extract_html_text,
        ".txt":  extract_plain_text,
        ".md":   extract_plain_text,
    }

    handler = handlers.get(extension)

    if handler:
        return handler(filepath)

    raise ValueError(
        f"Unsupported file type: {extension}. "
        f"Supported types: {', '.join(handlers.keys())}"
    )